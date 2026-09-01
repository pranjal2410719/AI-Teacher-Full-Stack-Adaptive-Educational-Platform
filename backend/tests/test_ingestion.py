"""
Comprehensive Test Suite for M1: Learning Material Ingestion & RAG Engine.
Tests multi-format parsers (PDF, DOCX, PPTX, TXT, MD), chunking, BM25 ranker,
NumpyVectorStore hybrid retrieval, topic parametric mode, unified LLM client, and REST APIs.
"""

import io
import pytest
from pathlib import Path
from fastapi.testclient import TestClient
import docx
import pptx
import pypdf
import numpy as np

from backend.app.config import settings
from backend.app.main import app
from backend.app.models.ingestion import (
    DocumentMetadata,
    DocumentChunk,
    TopicIngestionRequest,
    TopicIngestionResponse,
    RAGQuery,
    ChunkMatch,
    RAGResponse
)
from backend.app.services.llm_client import UnifiedLLMClient, llm_client
from backend.app.services.vector_store import (
    BM25Ranker,
    chunk_text_sliding_window,
    DocumentVectorIndex,
    NumpyVectorStore,
    vector_store
)
from backend.app.services.ingestion_service import (
    IngestionService,
    ingestion_service
)


@pytest.fixture
def client():
    return TestClient(app)


# -----------------------------------------------------------------------------
# 1. Pydantic Model Validation Tests
# -----------------------------------------------------------------------------

def test_models_validation():
    # Valid Topic Request
    req = TopicIngestionRequest(topic="Calculus Derivatives", subject_category="Mathematics", language="en")
    assert req.topic == "Calculus Derivatives"
    assert req.subject_category == "Mathematics"

    # Blank / whitespace topic rejection
    with pytest.raises(ValueError, match="cannot be blank"):
        TopicIngestionRequest(topic="   ")

    # Pure emoji / punctuation topic rejection
    with pytest.raises(ValueError, match="alphanumeric"):
        TopicIngestionRequest(topic="???!!! 🚀🎉")

    # Valid RAG Query
    q = RAGQuery(document_id="doc_123", query="What is the power rule?", top_k=5)
    assert q.top_k == 5
    assert q.query == "What is the power rule?"

    # Empty query rejection
    with pytest.raises(ValueError, match="cannot be empty"):
        RAGQuery(query="   ")


# -----------------------------------------------------------------------------
# 2. PDF Parser Tests
# -----------------------------------------------------------------------------

def test_pdf_parser_and_ingestion():
    # Create a genuine 2-page PDF in memory using pypdf
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=612, height=792)
    writer.add_blank_page(width=612, height=792)
    
    # We can write metadata and annotations or extract
    pdf_buffer = io.BytesIO()
    writer.write(pdf_buffer)
    pdf_bytes = pdf_buffer.getvalue()

    metadata, chunks = ingestion_service.ingest_document(pdf_bytes, "calculus_basics.pdf")
    assert metadata.document_id.startswith("doc_")
    assert metadata.file_type == "pdf"
    assert metadata.total_pages == 2
    assert len(chunks) > 0
    assert metadata.status == "ready"


def test_pdf_encrypted_handling():
    # Create an encrypted PDF
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=300, height=300)
    writer.encrypt("super_secret_password")
    
    pdf_buf = io.BytesIO()
    writer.write(pdf_buf)
    encrypted_bytes = pdf_buf.getvalue()

    with pytest.raises(ValueError, match="password-protected"):
        ingestion_service.ingest_document(encrypted_bytes, "locked.pdf")


# -----------------------------------------------------------------------------
# 3. DOCX Parser Tests
# -----------------------------------------------------------------------------

def test_docx_parser_and_ingestion():
    doc = docx.Document()
    doc.add_heading("Quantum Mechanics Fundamentals", level=1)
    doc.add_paragraph("Wave-particle duality posits that all particles exhibit both wave and particle properties.")
    doc.add_heading("Schrodinger Equation", level=2)
    doc.add_paragraph("The time-dependent Schrodinger equation describes how the quantum state of a physical system changes over time.")

    # Add a table
    table = doc.add_table(rows=2, cols=2)
    table.rows[0].cells[0].text = "Particle"
    table.rows[0].cells[1].text = "Spin"
    table.rows[1].cells[0].text = "Electron"
    table.rows[1].cells[1].text = "1/2"

    buf = io.BytesIO()
    doc.save(buf)
    docx_bytes = buf.getvalue()

    metadata, chunks = ingestion_service.ingest_document(docx_bytes, "quantum_physics.docx")
    assert metadata.file_type == "docx"
    assert len(chunks) >= 2
    assert any("Wave-particle duality" in c.text for c in chunks)
    assert any("Electron" in c.text and "Spin" in c.text for c in chunks)  # Table was parsed into markdown


def test_docx_xml_fallback_parser():
    # Test XML fallback directly
    doc = docx.Document()
    doc.add_paragraph("Fallback extraction test paragraph.")
    buf = io.BytesIO()
    doc.save(buf)
    raw_bytes = buf.getvalue()

    sections, all_text, count = ingestion_service._parse_docx_xml_fallback(raw_bytes)
    assert count > 0
    assert "Fallback extraction test paragraph." in all_text[0]


# -----------------------------------------------------------------------------
# 4. PPTX Parser Tests
# -----------------------------------------------------------------------------

def test_pptx_parser_and_ingestion():
    prs = pptx.Presentation()
    
    # Slide 1: Title and bullets
    s1 = prs.slides.add_slide(prs.slide_layouts[0])
    s1.shapes.title.text = "Data Structures: Trees and Graphs"
    s1.placeholders[1].text = "Binary Trees\nAVL Trees\nRed-Black Trees"
    s1.notes_slide.notes_text_frame.text = "Presenter note: Mention that balancing guarantees O(log N) operations."

    # Slide 2: Graph Theory
    s2 = prs.slides.add_slide(prs.slide_layouts[1])
    s2.shapes.title.text = "Graph Traversal Algorithms"
    s2.placeholders[1].text = "Breadth-First Search (BFS) uses a Queue.\nDepth-First Search (DFS) uses a Stack."

    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    metadata, chunks = ingestion_service.ingest_document(pptx_bytes, "trees_graphs.pptx")
    assert metadata.file_type == "pptx"
    assert metadata.total_pages == 2
    assert len(chunks) == 2
    # Verify speaker note inclusion
    assert any("guarantees O(log N)" in c.text for c in chunks)


def test_pptx_xml_fallback_parser():
    prs = pptx.Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = "Fallback Slide Title"
    buf = io.BytesIO()
    prs.save(buf)
    pptx_bytes = buf.getvalue()

    sections, all_text, count = ingestion_service._parse_pptx_xml_fallback(pptx_bytes)
    assert count == 1
    assert "Fallback Slide Title" in all_text[0]


# -----------------------------------------------------------------------------
# 5. TXT / Markdown Parser Tests
# -----------------------------------------------------------------------------

def test_txt_md_parser():
    md_content = """# Cell Biology

Cells are the basic structural, functional, and biological units of all known organisms.

## Mitochondria
The mitochondria is known as the powerhouse of the cell because it generates ATP.

## Nucleus
The nucleus contains the hereditary material (DNA) and coordinates cellular activities.
"""
    metadata, chunks = ingestion_service.ingest_document(md_content.encode("utf-8"), "biology.md")
    assert metadata.file_type == "md"
    assert len(chunks) >= 2
    assert any("powerhouse of the cell" in c.text for c in chunks)


def test_empty_txt_rejected():
    with pytest.raises(ValueError, match="empty"):
        ingestion_service.ingest_document(b"", "empty.txt")


# -----------------------------------------------------------------------------
# 6. Topic Parametric Mode Tests
# -----------------------------------------------------------------------------

def test_topic_parametric_ingestion():
    req = TopicIngestionRequest(
        topic="Binary Search Trees",
        subject_category="Computer Science",
        additional_notes="Focus on insertion and search complexity",
        language="en"
    )
    resp, chunks = ingestion_service.ingest_topic(req)
    assert resp.topic_id.startswith("top_")
    assert resp.topic == "Binary Search Trees"
    assert resp.subject_category == "Computer Science"
    assert len(chunks) >= 3
    assert resp.generated_chunks_count == len(chunks)

    # Verify RAG query works against parametric topic index
    rag_resp = ingestion_service.query_rag(
        RAGQuery(topic_id=resp.topic_id, query="What is the search time complexity of BST?", top_k=2)
    )
    assert rag_resp.total_results > 0
    assert rag_resp.results[0].similarity_score > 0
    assert len(rag_resp.grounded_context) > 0


# -----------------------------------------------------------------------------
# 7. BM25 Lexical Ranker Tests
# -----------------------------------------------------------------------------

def test_bm25_ranker():
    ranker = BM25Ranker(k1=1.5, b=0.75)
    docs = [
        "Photosynthesis in green plants converts water and carbon dioxide into oxygen and glucose.",
        "Newton's laws of motion describe the relationship between a body and the forces acting upon it.",
        "In computer programming, recursion is a method of solving problems where a function calls itself."
    ]
    ranker.fit(docs)

    # Query matching doc 0
    scores_bio = ranker.score("photosynthesis plants oxygen")
    assert scores_bio[0] > scores_bio[1]
    assert scores_bio[0] > scores_bio[2]
    assert scores_bio[0] == 1.0  # Normalized to max 1.0

    # Query matching doc 2
    scores_cs = ranker.score("recursion function calls itself")
    assert scores_cs[2] > scores_cs[0]
    assert scores_cs[2] > scores_cs[1]

    # Query with non-matching words
    scores_none = ranker.score("unmatched gibberish xyz123")
    assert all(s == 0.0 for s in scores_none)


# -----------------------------------------------------------------------------
# 8. NumpyVectorStore & Hybrid Search Tests
# -----------------------------------------------------------------------------

def test_numpy_vector_store_hybrid():
    chunks = [
        DocumentChunk(
            chunk_id="chk_calc_01",
            document_id="doc_calc",
            source_filename="calculus.pdf",
            page_or_slide=1,
            section_title="Limits",
            text="The formal definition of a limit involves epsilon and delta thresholds.",
            token_count=12
        ),
        DocumentChunk(
            chunk_id="chk_calc_02",
            document_id="doc_calc",
            source_filename="calculus.pdf",
            page_or_slide=2,
            section_title="Derivatives",
            text="The derivative represents the instantaneous rate of change and the slope of the tangent line.",
            token_count=16
        )
    ]
    
    idx = vector_store.add_document("doc_calc", chunks)
    assert idx.embeddings is not None
    assert idx.embeddings.shape == (2, settings.vector_dim)

    # Query for derivative slope
    rag_res = vector_store.query(
        query="tangent line slope instantaneous rate",
        target_id="doc_calc",
        top_k=1,
        alpha=0.6
    )
    assert rag_res.total_results == 1
    assert rag_res.results[0].chunk_id == "chk_calc_02"
    assert "tangent line" in rag_res.results[0].text
    assert rag_res.results[0].vector_score is not None
    assert rag_res.results[0].bm25_score is not None


# -----------------------------------------------------------------------------
# 9. Unified LLM Client Tests
# -----------------------------------------------------------------------------

def test_unified_llm_client_embeddings():
    texts = ["Differential Equations", "Calculus Integrals", "Ancient History Rome"]
    embeddings = llm_client.generate_embeddings(texts)
    assert len(embeddings) == 3
    assert len(embeddings[0]) == settings.vector_dim
    
    # Check L2 normalization
    norm = np.linalg.norm(np.array(embeddings[0]))
    assert np.isclose(norm, 1.0, atol=1e-4)


def test_unified_llm_client_json_extraction():
    raw_markdown_json = "Here is the response:\n```json\n{\"concept\": \"Derivatives\", \"order\": 1}\n```\nHope this helps!"
    parsed = UnifiedLLMClient.extract_json(raw_markdown_json)
    assert parsed["concept"] == "Derivatives"
    assert parsed["order"] == 1


# -----------------------------------------------------------------------------
# 10. REST API Endpoints Tests
# -----------------------------------------------------------------------------

def test_api_upload_docx(client):
    doc = docx.Document()
    doc.add_heading("API Upload Test", level=1)
    doc.add_paragraph("Testing FastAPI upload route with python-docx buffer.")
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)

    response = client.post(
        "/api/v1/materials/upload",
        files={"file": ("api_test.docx", buf.getvalue(), "application/vnd.openxmlformats-officedocument.wordprocessingml.document")}
    )
    assert response.status_code == 200
    data = response.json()
    assert data["filename"] == "api_test.docx"
    assert data["file_type"] == "docx"
    assert data["chunk_count"] > 0
    doc_id = data["document_id"]

    # Verify retrieval via GET /{doc_id}
    get_res = client.get(f"/api/v1/materials/{doc_id}")
    assert get_res.status_code == 200
    assert get_res.json()["document_id"] == doc_id


def test_api_upload_errors(client):
    # Empty file
    response = client.post(
        "/api/v1/materials/upload",
        files={"file": ("empty.txt", b"", "text/plain")}
    )
    assert response.status_code == 400
    assert "empty" in response.json()["detail"]

    # Unsupported extension
    response = client.post(
        "/api/v1/materials/upload",
        files={"file": ("program.exe", b"binary_data", "application/octet-stream")}
    )
    assert response.status_code == 400
    assert "Unsupported file extension" in response.json()["detail"]


def test_api_topic_ingest_and_query(client):
    # Ingest Topic
    resp = client.post(
        "/api/v1/materials/topic",
        json={"topic": "Linear Algebra Matrices", "subject_category": "Mathematics"}
    )
    assert resp.status_code == 200
    topic_data = resp.json()
    topic_id = topic_data["topic_id"]
    assert topic_data["topic"] == "Linear Algebra Matrices"

    # Query Topic
    q_resp = client.post(
        "/api/v1/materials/query",
        json={"topic_id": topic_id, "query": "matrix multiplication and determinants", "top_k": 3}
    )
    assert q_resp.status_code == 200
    q_data = q_resp.json()
    assert q_data["total_results"] > 0
    assert len(q_data["results"]) <= 3


def test_api_list_materials(client):
    resp = client.get("/api/v1/materials")
    assert resp.status_code == 200
    materials = resp.json()
    assert isinstance(materials, list)
    assert len(materials) > 0


def test_api_get_material_not_found(client):
    resp = client.get("/api/v1/materials/doc_non_existent_9999")
    assert resp.status_code == 404


# -----------------------------------------------------------------------------
# 11. Advanced Persistence & Boundary Tests
# -----------------------------------------------------------------------------

def test_persistence_reload_from_disk():
    # Ingest document
    meta, chunks = ingestion_service.ingest_document(
        b"Persistence verification content. Stored in disk vector indices.\n\nSecondary concept line.",
        "persistence_test.txt"
    )
    doc_id = meta.document_id

    # Create fresh vector store instance pointing to same directory
    new_store = NumpyVectorStore()
    assert doc_id not in new_store.indices

    # Query without pre-loading into memory (triggers load_from_disk)
    rag_res = new_store.query("Persistence verification", target_id=doc_id, top_k=2)
    assert rag_res.total_results > 0
    assert "Persistence verification" in rag_res.results[0].text
    assert doc_id in new_store.indices


def test_rag_query_exceeding_chunk_count():
    meta, chunks = ingestion_service.ingest_document(
        b"Single chunk test line.",
        "single.txt"
    )
    # Request top_k = 10 when only 1 chunk exists
    rag_res = ingestion_service.query_rag(RAGQuery(document_id=meta.document_id, query="single chunk", top_k=10))
    assert rag_res.total_results == 1
    assert len(rag_res.results) == 1


def test_multilingual_hindi_ingestion():
    hindi_text = """# प्रकाश संश्लेषण

प्रकाश संश्लेषण वह प्रक्रिया है जिसके द्वारा हरे पौधे सूर्य के प्रकाश की ऊर्जा को रासायनिक ऊर्जा में बदलते हैं।

## प्रकाशिक अभिक्रिया
क्लोरोफिल सूर्य के प्रकाश को अवशोषित करता है और एटीपी उत्पन्न करता है।
"""
    meta, chunks = ingestion_service.ingest_document(hindi_text.encode("utf-8"), "photosynthesis_hindi.md")
    assert meta.file_type == "md"
    assert len(chunks) >= 2

    # Query in Hindi
    rag_res = ingestion_service.query_rag(
        RAGQuery(document_id=meta.document_id, query="प्रकाश संश्लेषण और क्लोरोफिल", top_k=2)
    )
    assert rag_res.total_results > 0
    assert "प्रकाश संश्लेषण" in rag_res.grounded_context


def test_system_health_and_root_endpoints(client):
    res_root = client.get("/")
    assert res_root.status_code == 200
    assert "Welcome to AI Teacher" in res_root.json()["message"]

    res_health = client.get("/api/v1/health")
    assert res_health.status_code == 200
    health_data = res_health.json()
    assert health_data["status"] == "healthy"
    assert health_data["indexed_documents_count"] > 0
