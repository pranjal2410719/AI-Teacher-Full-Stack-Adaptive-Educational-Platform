"""
Document Ingestion and Parsing Service for AI Teacher.
Handles PDF, DOCX, PPT/PPTX, TXT/MD files, and Plain-Text Topic Parametric Generation.
"""

import os
import io
import re
import uuid
import json
import zipfile
import logging
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple, Union
from datetime import datetime, timezone
import xml.etree.ElementTree as ET

# Third-party parsers
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

from backend.app.config import settings
from backend.app.models.ingestion import (
    DocumentMetadata,
    DocumentChunk,
    TopicIngestionRequest,
    TopicIngestionResponse,
    RAGQuery,
    RAGResponse
)
from backend.app.services.llm_client import llm_client
from backend.app.services.vector_store import (
    vector_store,
    chunk_text_sliding_window
)

logger = logging.getLogger("ai_teacher.ingestion")


class ParsedSection:
    """Represents an extracted structural unit (page, slide, section, or chapter)."""
    def __init__(
        self,
        section_index: int,
        page_or_slide: Optional[int] = None,
        section_title: Optional[str] = None,
        text: str = "",
        extra_meta: Optional[Dict[str, Any]] = None
    ):
        self.section_index = section_index
        self.page_or_slide = page_or_slide
        self.section_title = section_title
        self.text = text
        self.extra_meta = extra_meta or {}


class ParsedDocument:
    """Aggregated parsed representation of an ingested document."""
    def __init__(
        self,
        filename: str,
        file_type: str,
        total_pages_or_slides: int,
        sections: List[ParsedSection],
        raw_full_text: str,
        summary: str = ""
    ):
        self.filename = filename
        self.file_type = file_type
        self.total_pages_or_slides = total_pages_or_slides
        self.sections = sections
        self.raw_full_text = raw_full_text
        self.summary = summary


class IngestionService:
    """
    Core service coordinating multi-format document parsing, semantic chunking,
    parametric topic synthesis, and vector indexing.
    """

    def __init__(self):
        self.upload_dir = settings.upload_dir
        self.metadata_registry: Dict[str, DocumentMetadata] = {}
        self._load_persisted_metadata()

    def _load_persisted_metadata(self) -> None:
        """Loads existing document metadata from disk on startup."""
        if not settings.vector_index_dir.exists():
            return
        for target_folder in settings.vector_index_dir.iterdir():
            if target_folder.is_dir():
                meta_file = target_folder / "metadata.json"
                if meta_file.exists():
                    try:
                        with open(meta_file, "r", encoding="utf-8") as f:
                            data = json.load(f)
                            doc_meta = DocumentMetadata(**data)
                            self.metadata_registry[doc_meta.document_id] = doc_meta
                    except Exception as e:
                        logger.warning(f"Could not load metadata from {meta_file}: {e}")

    def _persist_metadata(self, metadata: DocumentMetadata) -> None:
        """Saves document metadata to its respective vector index directory."""
        target_dir = settings.vector_index_dir / metadata.document_id
        target_dir.mkdir(parents=True, exist_ok=True)
        with open(target_dir / "metadata.json", "w", encoding="utf-8") as f:
            json.dump(metadata.model_dump(), f, indent=2)

    # -------------------------------------------------------------------------
    # Format Parsers
    # -------------------------------------------------------------------------

    def parse_pdf(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        """
        Extracts structured text per page from PDF using pypdf.
        Handles password-protected / encrypted PDFs gracefully.
        """
        if pypdf is None:
            raise RuntimeError("pypdf library is not installed.")

        stream = io.BytesIO(file_bytes)
        try:
            reader = pypdf.PdfReader(stream)
        except Exception as e:
            raise ValueError(f"Corrupted or unreadable PDF file: {e}")

        if reader.is_encrypted:
            # Try decrypting with empty string (standard for some encrypted PDFs)
            try:
                decrypted = reader.decrypt("")
                if decrypted == 0:
                    raise ValueError("PDF is password-protected. Please upload an unprotected file.")
            except Exception:
                raise ValueError("PDF is password-protected. Please upload an unprotected file.")

        total_pages = len(reader.pages)
        if total_pages == 0:
            raise ValueError("PDF contains 0 pages.")

        sections: List[ParsedSection] = []
        all_text_parts: List[str] = []

        for page_idx, page in enumerate(reader.pages, start=1):
            try:
                page_text = page.extract_text() or ""
            except Exception as e:
                logger.warning(f"Error extracting text from page {page_idx} of {filename}: {e}")
                page_text = ""

            page_text = page_text.strip()
            if not page_text:
                continue

            # Detect potential heading from first non-empty line
            lines = [l.strip() for l in page_text.splitlines() if l.strip()]
            section_title = lines[0][:80] if lines else f"Page {page_idx}"

            sections.append(
                ParsedSection(
                    section_index=len(sections),
                    page_or_slide=page_idx,
                    section_title=section_title,
                    text=page_text,
                    extra_meta={"page_number": page_idx}
                )
            )
            all_text_parts.append(f"--- Page {page_idx} ---\n{page_text}")

        raw_full_text = "\n\n".join(all_text_parts)
        if not raw_full_text.strip():
            # If scan/image-only PDF
            raw_full_text = f"Scanned or image-only PDF: {filename}. (No selectable text extracted)."
            sections.append(
                ParsedSection(
                    section_index=0,
                    page_or_slide=1,
                    section_title="Document Content",
                    text=raw_full_text
                )
            )

        summary = self._generate_summary_from_text(raw_full_text, filename)

        return ParsedDocument(
            filename=filename,
            file_type="pdf",
            total_pages_or_slides=total_pages,
            sections=sections,
            raw_full_text=raw_full_text,
            summary=summary
        )

    def parse_docx(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        """
        Extracts paragraphs, heading hierarchy, and tables from DOCX.
        Includes XML fallback if python-docx fails.
        """
        sections: List[ParsedSection] = []
        all_text_parts: List[str] = []
        total_elements = 0

        if docx is not None:
            try:
                stream = io.BytesIO(file_bytes)
                doc = docx.Document(stream)
                current_heading = "Introduction / Overview"
                current_paragraphs: List[str] = []
                section_count = 0

                for para in doc.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    total_elements += 1

                    # Check if paragraph is a heading
                    style_name = para.style.name.lower() if para.style else ""
                    if "heading" in style_name or "title" in style_name:
                        if current_paragraphs:
                            body = "\n\n".join(current_paragraphs)
                            sections.append(
                                ParsedSection(
                                    section_index=section_count,
                                    page_or_slide=section_count + 1,
                                    section_title=current_heading,
                                    text=body
                                )
                            )
                            all_text_parts.append(f"### {current_heading}\n{body}")
                            section_count += 1
                            current_paragraphs = []
                        current_heading = text
                    else:
                        current_paragraphs.append(text)

                # Process Tables into Markdown formatted strings
                for table in doc.tables:
                    table_rows = []
                    for row in table.rows:
                        cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                        table_rows.append(cells)

                    if table_rows:
                        # Convert to Markdown table
                        header = table_rows[0]
                        md_table = "| " + " | ".join(header) + " |\n"
                        md_table += "| " + " | ".join(["---"] * len(header)) + " |\n"
                        for row_data in table_rows[1:]:
                            # Pad row if mismatched
                            padded = row_data + [""] * (len(header) - len(row_data))
                            md_table += "| " + " | ".join(padded[:len(header)]) + " |\n"

                        current_paragraphs.append(f"Table Data:\n{md_table}")
                        total_elements += 1

                if current_paragraphs:
                    body = "\n\n".join(current_paragraphs)
                    sections.append(
                        ParsedSection(
                            section_index=section_count,
                            page_or_slide=section_count + 1,
                            section_title=current_heading,
                            text=body
                        )
                    )
                    all_text_parts.append(f"### {current_heading}\n{body}")

            except Exception as e:
                logger.warning(f"python-docx parsing failed on {filename}: {e}. Trying XML fallback.")
                sections, all_text_parts, total_elements = self._parse_docx_xml_fallback(file_bytes)
        else:
            sections, all_text_parts, total_elements = self._parse_docx_xml_fallback(file_bytes)

        raw_full_text = "\n\n".join(all_text_parts)
        if not raw_full_text.strip():
            raw_full_text = f"Empty DOCX document: {filename}."

        summary = self._generate_summary_from_text(raw_full_text, filename)

        return ParsedDocument(
            filename=filename,
            file_type="docx",
            total_pages_or_slides=max(1, len(sections)),
            sections=sections,
            raw_full_text=raw_full_text,
            summary=summary
        )

    def _parse_docx_xml_fallback(self, file_bytes: bytes) -> Tuple[List[ParsedSection], List[str], int]:
        """Fallback DOCX extractor using raw word/document.xml extraction."""
        sections = []
        all_text = []
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                xml_content = z.read("word/document.xml")
                root = ET.fromstring(xml_content)
                namespaces = {"w": "http://schemas.openxmlformats.org/wordprocessingml/2006/main"}
                
                paras = []
                for p in root.findall(".//w:p", namespaces):
                    p_text = "".join(t.text for t in p.findall(".//w:t", namespaces) if t.text)
                    if p_text.strip():
                        paras.append(p_text.strip())

                if paras:
                    combined = "\n\n".join(paras)
                    sections.append(
                        ParsedSection(
                            section_index=0,
                            page_or_slide=1,
                            section_title="Document Content",
                            text=combined
                        )
                    )
                    all_text.append(combined)
                    return sections, all_text, len(paras)
        except Exception as e:
            logger.error(f"DOCX XML fallback also failed: {e}")

        return sections, all_text, 0

    def parse_pptx(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        """
        Extracts slide titles, text frames, bullet points, table contents,
        and presenter speaker notes from PPTX.
        """
        sections: List[ParsedSection] = []
        all_text_parts: List[str] = []
        total_slides = 0

        if pptx is not None:
            try:
                stream = io.BytesIO(file_bytes)
                prs = pptx.Presentation(stream)
                total_slides = len(prs.slides)

                for slide_idx, slide in enumerate(prs.slides, start=1):
                    slide_title = f"Slide {slide_idx}"
                    # Try to extract slide title shape
                    if slide.shapes.title and slide.shapes.title.text:
                        slide_title = slide.shapes.title.text.strip()

                    slide_text_lines: List[str] = []
                    # Extract from all shapes
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            for paragraph in shape.text_frame.paragraphs:
                                t = paragraph.text.strip()
                                if t and t != slide_title:
                                    slide_text_lines.append(t)
                        elif shape.has_table:
                            for row in shape.table.rows:
                                row_cells = [c.text.strip() for c in row.cells if c.text.strip()]
                                if row_cells:
                                    slide_text_lines.append(" | ".join(row_cells))

                    # Extract presenter speaker notes
                    speaker_notes = ""
                    if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                        notes_raw = slide.notes_slide.notes_text_frame.text.strip()
                        if notes_raw:
                            speaker_notes = f"\n[Speaker Notes: {notes_raw}]"

                    body = "\n".join(slide_text_lines)
                    full_slide_content = f"Title: {slide_title}\n{body}{speaker_notes}".strip()

                    sections.append(
                        ParsedSection(
                            section_index=slide_idx - 1,
                            page_or_slide=slide_idx,
                            section_title=slide_title,
                            text=full_slide_content,
                            extra_meta={"slide_number": slide_idx, "speaker_notes": bool(speaker_notes)}
                        )
                    )
                    all_text_parts.append(f"--- Slide {slide_idx}: {slide_title} ---\n{body}{speaker_notes}")

            except Exception as e:
                logger.warning(f"python-pptx parsing failed on {filename}: {e}. Trying PPTX XML fallback.")
                sections, all_text_parts, total_slides = self._parse_pptx_xml_fallback(file_bytes)
        else:
            sections, all_text_parts, total_slides = self._parse_pptx_xml_fallback(file_bytes)

        raw_full_text = "\n\n".join(all_text_parts)
        if not raw_full_text.strip():
            raw_full_text = f"Empty presentation: {filename}."

        summary = self._generate_summary_from_text(raw_full_text, filename)

        return ParsedDocument(
            filename=filename,
            file_type="pptx",
            total_pages_or_slides=max(1, total_slides),
            sections=sections,
            raw_full_text=raw_full_text,
            summary=summary
        )

    def _parse_pptx_xml_fallback(self, file_bytes: bytes) -> Tuple[List[ParsedSection], List[str], int]:
        """Fallback PPTX extractor using ppt/slides/slide*.xml parsing."""
        sections = []
        all_text = []
        slide_count = 0
        try:
            with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
                slide_files = sorted([name for name in z.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])
                slide_count = len(slide_files)

                for idx, s_name in enumerate(slide_files, start=1):
                    xml_content = z.read(s_name)
                    root = ET.fromstring(xml_content)
                    text_parts = [elem.text.strip() for elem in root.iter() if elem.tag.endswith("}t") and elem.text and elem.text.strip()]
                    
                    slide_title = text_parts[0] if text_parts else f"Slide {idx}"
                    body = "\n".join(text_parts[1:]) if len(text_parts) > 1 else ""
                    content = f"Title: {slide_title}\n{body}".strip()

                    sections.append(
                        ParsedSection(
                            section_index=idx - 1,
                            page_or_slide=idx,
                            section_title=slide_title,
                            text=content
                        )
                    )
                    all_text.append(f"--- Slide {idx}: {slide_title} ---\n{body}")
        except Exception as e:
            logger.error(f"PPTX XML fallback also failed: {e}")

        return sections, all_text, slide_count

    def parse_txt_md(self, file_bytes: bytes, filename: str) -> ParsedDocument:
        """
        Parses TXT and Markdown files with multi-encoding detection and Markdown header splitting.
        """
        encodings = ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]
        raw_text = None

        for enc in encodings:
            try:
                raw_text = file_bytes.decode(enc)
                break
            except (UnicodeDecodeError, LookupError):
                continue

        if raw_text is None:
            raw_text = file_bytes.decode("utf-8", errors="replace")

        raw_text = raw_text.strip()
        if not raw_text:
            raise ValueError(f"Text file {filename} is empty.")

        sections: List[ParsedSection] = []
        all_text_parts: List[str] = []

        # Split along markdown headers `# `, `## `, `### `
        header_splits = re.split(r"(?m)^(#{1,3}\s+.+)$", raw_text)

        if len(header_splits) > 1:
            current_title = "Introduction"
            current_body = []
            sec_idx = 0

            for part in header_splits:
                part = part.strip()
                if not part:
                    continue
                if part.startswith("#"):
                    if current_body:
                        text_val = "\n\n".join(current_body)
                        sections.append(
                            ParsedSection(
                                section_index=sec_idx,
                                page_or_slide=sec_idx + 1,
                                section_title=current_title,
                                text=text_val
                            )
                        )
                        all_text_parts.append(f"{current_title}\n{text_val}")
                        sec_idx += 1
                        current_body = []
                    current_title = part.lstrip("#").strip()
                else:
                    current_body.append(part)

            if current_body:
                text_val = "\n\n".join(current_body)
                sections.append(
                    ParsedSection(
                        section_index=sec_idx,
                        page_or_slide=sec_idx + 1,
                        section_title=current_title,
                        text=text_val
                    )
                )
                all_text_parts.append(f"{current_title}\n{text_val}")
        else:
            # Single or paragraph delimited text
            paragraphs = [p.strip() for p in raw_text.split("\n\n") if p.strip()]
            for idx, p in enumerate(paragraphs):
                title = p.splitlines()[0][:60] if p.splitlines() else f"Section {idx+1}"
                sections.append(
                    ParsedSection(
                        section_index=idx,
                        page_or_slide=idx + 1,
                        section_title=title,
                        text=p
                    )
                )
            all_text_parts.append(raw_text)

        summary = self._generate_summary_from_text(raw_text, filename)

        return ParsedDocument(
            filename=filename,
            file_type="md" if filename.lower().endswith(".md") else "txt",
            total_pages_or_slides=max(1, len(sections)),
            sections=sections,
            raw_full_text=raw_full_text if (raw_full_text := "\n\n".join(all_text_parts)) else raw_text,
            summary=summary
        )

    # -------------------------------------------------------------------------
    # Helper & Summary Generators
    # -------------------------------------------------------------------------

    def _generate_summary_from_text(self, text: str, filename: str) -> str:
        """Synthesizes high-level summary of text for DocumentMetadata."""
        words = text.split()
        sample = " ".join(words[:400])
        if len(words) < 30:
            return f"Summary of {filename}: {text[:200]}"

        prompt = (
            f"Please provide a concise 2-3 sentence educational overview summarizing the main topics "
            f"covered in the following document ({filename}):\n\n{sample}\n\nSummary:"
        )
        try:
            summary = llm_client.generate_completion(prompt, max_tokens=150, temperature=0.3)
            return summary.strip()
        except Exception:
            return f"Comprehensive educational material on {filename}, covering foundational principles and examples."

    # -------------------------------------------------------------------------
    # Main Ingestion Pipelines
    # -------------------------------------------------------------------------

    def ingest_document(
        self,
        file_bytes: bytes,
        filename: str,
        content_type: Optional[str] = None
    ) -> Tuple[DocumentMetadata, List[DocumentChunk]]:
        """
        Complete end-to-end ingestion pipeline:
        1. Validate file size and type.
        2. Parse file into structured ParsedDocument.
        3. Perform structure-aware chunking.
        4. Embed and index chunks into NumpyVectorStore with BM25 ranker.
        5. Persist uploaded file and metadata.
        """
        if not file_bytes or len(file_bytes) == 0:
            raise ValueError("File is empty (0 bytes). Please upload valid educational material.")

        max_bytes = settings.max_upload_size_mb * 1024 * 1024
        if len(file_bytes) > max_bytes:
            raise ValueError(f"File size ({len(file_bytes)/(1024*1024):.1f}MB) exceeds maximum limit of {settings.max_upload_size_mb}MB.")

        ext = Path(filename).suffix.lower().lstrip(".")
        if ext not in ["pdf", "docx", "pptx", "ppt", "txt", "md"]:
            raise ValueError(f"Unsupported file extension '.{ext}'. Supported formats: PDF, DOCX, PPT, PPTX, TXT, MD.")

        doc_id = f"doc_{uuid.uuid4().hex[:10]}"

        # Dispatch to format-specific parser
        if ext == "pdf":
            parsed = self.parse_pdf(file_bytes, filename)
        elif ext == "docx":
            parsed = self.parse_docx(file_bytes, filename)
        elif ext in ["pptx", "ppt"]:
            parsed = self.parse_pptx(file_bytes, filename)
        else:
            parsed = self.parse_txt_md(file_bytes, filename)

        # Structure-aware chunking
        all_chunks: List[DocumentChunk] = []
        chunk_counter = 0

        for section in parsed.sections:
            sec_chunks = chunk_text_sliding_window(
                text=section.text,
                chunk_size=settings.default_chunk_size,
                overlap=settings.default_chunk_overlap,
                source_filename=filename,
                document_id=doc_id,
                page_or_slide=section.page_or_slide,
                section_title=section.section_title,
                start_chunk_index=chunk_counter
            )
            all_chunks.extend(sec_chunks)
            chunk_counter += len(sec_chunks)

        if not all_chunks:
            # Fallback single chunk
            all_chunks.append(
                DocumentChunk(
                    chunk_id=f"chk_{doc_id}_0000",
                    document_id=doc_id,
                    source_filename=filename,
                    page_or_slide=1,
                    section_title="Overview",
                    text=parsed.raw_full_text[:1000],
                    token_count=len(parsed.raw_full_text[:1000].split()),
                    chunk_index=0
                )
            )

        # Save raw file to upload storage
        safe_filename = re.sub(r"[^a-zA-Z0-9_.-]", "_", filename)
        saved_file_path = self.upload_dir / f"{doc_id}_{safe_filename}"
        with open(saved_file_path, "wb") as f:
            f.write(file_bytes)

        # Index in Vector Store
        vector_store.add_document(target_id=doc_id, chunks=all_chunks)

        # Create Metadata
        metadata = DocumentMetadata(
            document_id=doc_id,
            filename=filename,
            file_type=ext,
            file_size_bytes=len(file_bytes),
            total_pages=parsed.total_pages_or_slides,
            chunk_count=len(all_chunks),
            extracted_summary=parsed.summary,
            created_at=datetime.now(timezone.utc).isoformat(),
            status="ready",
            metadata_extra={
                "stored_path": str(saved_file_path),
                "total_sections": len(parsed.sections)
            }
        )

        self.metadata_registry[doc_id] = metadata
        self._persist_metadata(metadata)

        logger.info(f"Ingested document {doc_id} ('{filename}'): {len(all_chunks)} chunks indexed.")
        return metadata, all_chunks

    def ingest_topic(self, request: TopicIngestionRequest) -> Tuple[TopicIngestionResponse, List[DocumentChunk]]:
        """
        Parametric knowledge ingestion when no file is uploaded:
        1. Generates structured seed syllabus via LLM client or parametric generator.
        2. Chunks the concepts, definitions, step-by-step walkthrough, and misconceptions.
        3. Indexes chunks into NumpyVectorStore for grounded lesson planning.
        """
        topic_id = f"top_{uuid.uuid4().hex[:10]}"
        topic_clean = request.topic.strip()
        category = request.subject_category or "General"

        # Generate structured parametric seed syllabus
        prompt = (
            f"Generate structured educational seed syllabus for the topic: '{topic_clean}' in subject area: '{category}'.\n"
            f"Target language: {request.language or 'en'}.\n"
            f"Additional guidelines: {request.additional_notes or 'Comprehensive overview'}.\n"
            f"Provide output in structured JSON format with fields: 'topic', 'summary', 'core_concepts' "
            f"(with 'concept_title', 'description', 'key_points'), and 'sample_questions' (with 'question', 'answer')."
        )

        raw_output = llm_client.generate_completion(
            prompt=prompt,
            system_prompt="You are an expert pedagogical curriculum designer.",
            json_mode=True,
            temperature=0.5
        )

        try:
            syllabus_data = llm_client.extract_json(raw_output)
        except Exception:
            syllabus_data = {
                "topic": topic_clean,
                "summary": f"Comprehensive educational foundation covering core mechanics and examples for {topic_clean}.",
                "core_concepts": [
                    {
                        "concept_title": f"Foundations of {topic_clean}",
                        "description": f"Core definitions and key principles of {topic_clean}.",
                        "key_points": ["Basic definitions", "Essential notation", "Fundamental rules"]
                    },
                    {
                        "concept_title": f"Worked Examples and Practice in {topic_clean}",
                        "description": f"Applied demonstrations and step-by-step solutions for {topic_clean}.",
                        "key_points": ["Step-by-step derivation", "Real-world application", "Self-check verification"]
                    }
                ],
                "sample_questions": [
                    {
                        "question": f"What is the central concept behind {topic_clean}?",
                        "answer": f"The fundamental rule and primary mechanism of {topic_clean}."
                    }
                ]
            }

        seed_summary = syllabus_data.get("summary", f"Structured curriculum for {topic_clean}.")

        # Convert syllabus components into DocumentChunk objects
        chunks: List[DocumentChunk] = []
        chunk_idx = 0

        # Chunk 1: Overview and Summary
        chunks.append(
            DocumentChunk(
                chunk_id=f"chk_{topic_id}_{chunk_idx:04d}",
                document_id=topic_id,
                source_filename=f"Topic: {topic_clean}",
                page_or_slide=1,
                section_title="Topic Overview & Objectives",
                text=f"Topic: {topic_clean}\nCategory: {category}\n\nSummary:\n{seed_summary}",
                token_count=len(seed_summary.split()) + 20,
                chunk_index=chunk_idx
            )
        )
        chunk_idx += 1

        # Chunks for each core concept
        for i, concept in enumerate(syllabus_data.get("core_concepts", []), start=1):
            title = concept.get("concept_title", f"Concept {i}")
            desc = concept.get("description", "")
            points = "\n- ".join(concept.get("key_points", []))
            concept_text = f"Concept: {title}\n\nDescription:\n{desc}\n\nKey Principles:\n- {points}"

            chunks.append(
                DocumentChunk(
                    chunk_id=f"chk_{topic_id}_{chunk_idx:04d}",
                    document_id=topic_id,
                    source_filename=f"Topic: {topic_clean}",
                    page_or_slide=i + 1,
                    section_title=title,
                    text=concept_text,
                    token_count=len(concept_text.split()),
                    chunk_index=chunk_idx
                )
            )
            chunk_idx += 1

        # Chunks for questions and misconceptions
        questions = syllabus_data.get("sample_questions", [])
        if questions:
            q_lines = []
            for q in questions:
                q_text = q.get("question", "")
                a_text = q.get("answer", "")
                q_lines.append(f"Q: {q_text}\nA: {a_text}")
            qa_block = "Diagnostic Questions & Review:\n\n" + "\n\n".join(q_lines)

            chunks.append(
                DocumentChunk(
                    chunk_id=f"chk_{topic_id}_{chunk_idx:04d}",
                    document_id=topic_id,
                    source_filename=f"Topic: {topic_clean}",
                    page_or_slide=len(syllabus_data.get("core_concepts", [])) + 2,
                    section_title="Diagnostic Checks & Misconceptions",
                    text=qa_block,
                    token_count=len(qa_block.split()),
                    chunk_index=chunk_idx
                )
            )
            chunk_idx += 1

        # Index in Vector Store
        vector_store.add_document(target_id=topic_id, chunks=chunks)

        # Create Metadata
        metadata = DocumentMetadata(
            document_id=topic_id,
            filename=f"Topic: {topic_clean}",
            file_type="topic",
            file_size_bytes=sum(len(c.text) for c in chunks),
            total_pages=len(chunks),
            chunk_count=len(chunks),
            extracted_summary=seed_summary,
            created_at=datetime.now(timezone.utc).isoformat(),
            status="ready",
            metadata_extra={
                "topic": topic_clean,
                "subject_category": category,
                "language": request.language or "en"
            }
        )

        self.metadata_registry[topic_id] = metadata
        self._persist_metadata(metadata)

        response = TopicIngestionResponse(
            topic_id=topic_id,
            topic=topic_clean,
            subject_category=category,
            seed_summary=seed_summary,
            generated_chunks_count=len(chunks),
            created_at=metadata.created_at,
            status="ready"
        )

        logger.info(f"Ingested topic {topic_id} ('{topic_clean}'): {len(chunks)} parametric chunks indexed.")
        return response, chunks

    def query_rag(self, query: RAGQuery) -> RAGResponse:
        """Executes grounded RAG retrieval across document or topic index."""
        target_id = query.document_id or query.topic_id
        return vector_store.query(
            query=query.query,
            target_id=target_id,
            top_k=query.top_k,
            alpha=settings.hybrid_alpha
        )

    def get_metadata(self, document_id: str) -> Optional[DocumentMetadata]:
        """Retrieves document or topic metadata by ID."""
        if document_id in self.metadata_registry:
            return self.metadata_registry[document_id]
        # Check disk
        meta_file = settings.vector_index_dir / document_id / "metadata.json"
        if meta_file.exists():
            with open(meta_file, "r", encoding="utf-8") as f:
                data = json.load(f)
                meta = DocumentMetadata(**data)
                self.metadata_registry[document_id] = meta
                return meta
        return None

    def list_all_materials(self) -> List[DocumentMetadata]:
        """Lists all registered documents and topics."""
        return list(self.metadata_registry.values())


# Global shared singleton
ingestion_service = IngestionService()
