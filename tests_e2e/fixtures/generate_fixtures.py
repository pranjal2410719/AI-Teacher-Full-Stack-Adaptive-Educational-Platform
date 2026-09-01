"""
Fixture generator script for AI Teacher 4-Tier E2E Test Suite.
Generates authentic PDF, DOCX, PPTX, and TXT files for Math, CS, Biology, and History.
"""

import os
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
import docx
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

FIXTURES_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__)))
os.makedirs(FIXTURES_DIR, exist_ok=True)

def generate_pdf_calculus():
    pdf_path = os.path.join(FIXTURES_DIR, "calculus_limits.pdf")
    doc = SimpleDocTemplate(pdf_path, pagesize=letter, rightMargin=40, leftMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    story = []
    
    title_style = ParagraphStyle(
        'DocTitle',
        parent=styles['Heading1'],
        fontSize=20,
        leading=24,
        textColor=colors.HexColor("#1e293b"),
        spaceAfter=12
    )
    h2_style = ParagraphStyle(
        'DocH2',
        parent=styles['Heading2'],
        fontSize=14,
        leading=18,
        textColor=colors.HexColor("#2563eb"),
        spaceBefore=10,
        spaceAfter=6
    )
    body_style = ParagraphStyle(
        'DocBody',
        parent=styles['Normal'],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#334155"),
        spaceAfter=8
    )
    formula_style = ParagraphStyle(
        'DocFormula',
        parent=styles['Code'],
        fontSize=10,
        leading=14,
        textColor=colors.HexColor("#0f172a"),
        backColor=colors.HexColor("#f1f5f9"),
        borderColor=colors.HexColor("#cbd5e1"),
        borderWidth=1,
        borderPadding=6,
        spaceAfter=10
    )

    story.append(Paragraph("Chapter 1: Limits, Continuity, and the Foundations of Calculus", title_style))
    story.append(Paragraph("<b>Course</b>: Advanced Mathematics &amp; Calculus Foundations | <b>Target Level</b>: High School / Undergraduate", body_style))
    story.append(Spacer(1, 10))

    story.append(Paragraph("1.1 Intuitive Concept of a Limit", h2_style))
    story.append(Paragraph("In calculus, the limit is the foundational concept upon which differentiation and integration rest. Intuitively, when we say the limit of f(x) as x approaches c equals L, written as lim_{x->c} f(x) = L, we mean that we can make the value of f(x) arbitrarily close to L by taking x sufficiently close to c, but not equal to c.", body_style))
    story.append(Paragraph("Consider the function f(x) = (x^2 - 1) / (x - 1). At x = 1, f(1) is undefined (0/0 indeterminate form). However, for x != 1, factoring gives f(x) = (x-1)(x+1)/(x-1) = x+1. As x approaches 1 from either side, f(x) approaches 2. Thus, lim_{x->1} (x^2 - 1)/(x - 1) = 2.", body_style))

    story.append(Paragraph("1.2 The Formal Epsilon-Delta Definition", h2_style))
    story.append(Paragraph("The rigorous Cauchy-Weierstrass definition states that lim_{x->c} f(x) = L if and only if:", body_style))
    story.append(Paragraph("For every epsilon > 0, there exists a delta > 0 such that whenever 0 &lt; |x - c| &lt; delta, it follows that |f(x) - L| &lt; epsilon.", formula_style))
    story.append(Paragraph("Here, epsilon represents the maximum allowable error bound in the output space, while delta represents the corresponding required proximity in the input domain.", body_style))

    story.append(Paragraph("1.3 One-Sided Limits and Continuity", h2_style))
    story.append(Paragraph("A two-sided limit exists if and only if the left-hand limit and right-hand limit both exist and are strictly equal:", body_style))
    story.append(Paragraph("lim_{x->c} f(x) = L  &lt;==&gt;  lim_{x->c^-} f(x) = L  and  lim_{x->c^+} f(x) = L", formula_style))
    story.append(Paragraph("A function f(x) is continuous at x = c if: (1) f(c) is defined, (2) lim_{x->c} f(x) exists, and (3) lim_{x->c} f(x) = f(c).", body_style))

    story.append(Paragraph("1.4 The Derivative as a Limit of the Difference Quotient", h2_style))
    story.append(Paragraph("The instantaneous rate of change (derivative) of f(x) at x = a is defined as the limit of secant line slopes:", body_style))
    story.append(Paragraph("f'(a) = lim_{h->0} [f(a + h) - f(a)] / h", formula_style))
    story.append(Paragraph("Geometrically, as h approaches 0, the secant line passing through (a, f(a)) and (a+h, f(a+h)) rotates into the tangent line at x = a.", body_style))

    doc.build(story)
    print(f"Generated PDF: {pdf_path} ({os.path.getsize(pdf_path)} bytes)")

def generate_docx_bst():
    docx_path = os.path.join(FIXTURES_DIR, "binary_search_trees.docx")
    doc = Document()
    
    doc.add_heading("Data Structures: Binary Search Trees (BST)", level=0)
    
    p = doc.add_paragraph()
    p.add_run("Subject: ").bold = True
    p.add_run("Computer Science / Algorithms | ")
    p.add_run("Target Level: ").bold = True
    p.add_run("College / Intermediate")

    doc.add_heading("1. Definition and Binary Search Tree Invariant", level=1)
    doc.add_paragraph(
        "A Binary Search Tree (BST) is a hierarchical binary tree data structure where each node contains a key, "
        "and satisfies the binary search invariant: for any node N, all keys in N's left subtree are strictly less than N.key, "
        "and all keys in N's right subtree are strictly greater than N.key (assuming unique keys)."
    )

    doc.add_heading("2. Node Structure in Python", level=1)
    code_text = (
        "class TreeNode:\n"
        "    def __init__(self, val=0, left=None, right=None):\n"
        "        self.val = val\n"
        "        self.left = left\n"
        "        self.right = right\n"
    )
    doc.add_paragraph(code_text, style='Quote')

    doc.add_heading("3. Core Operations: Insertion and Search", level=1)
    doc.add_paragraph(
        "Recursive Insertion: To insert value x into tree rooted at root:\n"
        "1. If root is None, return TreeNode(x).\n"
        "2. If x < root.val, recursively insert into root.left: root.left = insert(root.left, x).\n"
        "3. If x > root.val, recursively insert into root.right: root.right = insert(root.right, x).\n"
        "4. Return root."
    )
    doc.add_paragraph(
        "Lookup / Search: Starting at root, compare target key with current node value. "
        "If equal, node is found. If target < root.val, traverse left; if target > root.val, traverse right. "
        "If null pointer is reached, key does not exist."
    )

    doc.add_heading("4. Time and Space Complexity Analysis", level=1)
    table = doc.add_table(rows=1, cols=3)
    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Operation"
    hdr_cells[1].text = "Average Case (Balanced)"
    hdr_cells[2].text = "Worst Case (Degenerate / Skewed)"

    ops_data = [
        ("Search", "O(log n)", "O(n) - Skewed linked list"),
        ("Insertion", "O(log n)", "O(n) - Skewed linked list"),
        ("Deletion", "O(log n)", "O(n) - Skewed linked list"),
        ("In-Order Traversal", "O(n)", "O(n) - Sorted order output"),
    ]
    for op, avg, worst in ops_data:
        row_cells = table.add_row().cells
        row_cells[0].text = op
        row_cells[1].text = avg
        row_cells[2].text = worst

    doc.add_paragraph()
    doc.add_paragraph(
        "Note on Self-Balancing Trees: To prevent O(n) degeneration when inserting sorted data, self-balancing BSTs "
        "such as AVL Trees or Red-Black Trees maintain a balance factor using tree rotations, guaranteeing O(log n) worst-case height."
    )

    doc.save(docx_path)
    print(f"Generated DOCX: {docx_path} ({os.path.getsize(docx_path)} bytes)")

def generate_pptx_biology():
    pptx_path = os.path.join(FIXTURES_DIR, "cell_biology.pptx")
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(5.625) # 16:9 widescreen
    blank_layout = prs.slide_layouts[6]

    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(blank_layout)
    txBox = slide1.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(2.5))
    tf = txBox.text_frame
    p1 = tf.paragraphs[0]
    p1.text = "Cell Biology: Structure and Function of Eukaryotic Organelles"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = RGBColor(30, 41, 59)
    p2 = tf.add_paragraph()
    p2.text = "High School & AP Biology | Comprehensive Organelle Breakdown"
    p2.font.size = Pt(16)
    p2.font.color.rgb = RGBColor(71, 85, 105)

    # Slide 2: Plasma Membrane & Nucleus
    slide2 = prs.slides.add_slide(blank_layout)
    txBox2 = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(4.5))
    tf2 = txBox2.text_frame
    p = tf2.paragraphs[0]
    p.text = "1. Plasma Membrane and the Nucleus"
    p.font.size = Pt(22)
    p.font.bold = True
    
    p = tf2.add_paragraph()
    p.text = "• Plasma Membrane: Phospholipid bilayer with embedded transport proteins, cholesterol, and glycoproteins following the fluid mosaic model. Regulates selective permeability (homeostasis)."
    p.font.size = Pt(14)
    p = tf2.add_paragraph()
    p.text = "• Nucleus: Houses genomic DNA organized into chromatin. Enclosed by double nuclear membrane with nuclear pores."
    p.font.size = Pt(14)
    p = tf2.add_paragraph()
    p.text = "• Nucleolus: Dense subnuclear region responsible for ribosomal RNA (rRNA) transcription and ribosome subunit assembly."
    p.font.size = Pt(14)

    # Slide 3: Mitochondria & Energy Conversion
    slide3 = prs.slides.add_slide(blank_layout)
    txBox3 = slide3.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(4.5))
    tf3 = txBox3.text_frame
    p = tf3.paragraphs[0]
    p.text = "2. Mitochondria: Cellular Respiration and ATP"
    p.font.size = Pt(22)
    p.font.bold = True
    
    p = tf3.add_paragraph()
    p.text = "• Structure: Outer membrane, intermembrane space, highly folded inner membrane (cristae), and mitochondrial matrix."
    p.font.size = Pt(14)
    p = tf3.add_paragraph()
    p.text = "• Function: Site of Krebs cycle (matrix) and Oxidative Phosphorylation / Electron Transport Chain (cristae)."
    p.font.size = Pt(14)
    p = tf3.add_paragraph()
    p.text = "• Endosymbiotic Theory: Mitochondria contain circular DNA and 70S ribosomes, indicating ancestral bacterial origin."
    p.font.size = Pt(14)
    p = tf3.add_paragraph()
    p.text = "• ATP Output: Generates ~30-32 ATP molecules per glucose molecule through aerobic respiration."
    p.font.size = Pt(14)

    # Slide 4: Plant vs Animal Cells
    slide4 = prs.slides.add_slide(blank_layout)
    txBox4 = slide4.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(8.4), Inches(4.5))
    tf4 = txBox4.text_frame
    p = tf4.paragraphs[0]
    p.text = "3. Distinct Structures: Plant vs Animal Cells"
    p.font.size = Pt(22)
    p.font.bold = True
    
    p = tf4.add_paragraph()
    p.text = "• Chloroplasts: Thylakoid stacks (grana) containing chlorophyll for photosynthesis (light & dark reactions) in plant cells."
    p.font.size = Pt(14)
    p = tf4.add_paragraph()
    p.text = "• Cell Wall: Rigid outer cellulose matrix in plant cells providing turgor pressure and structural support."
    p.font.size = Pt(14)
    p = tf4.add_paragraph()
    p.text = "• Central Vacuole: Large plant organelle for water storage, osmotic balance, and waste degradation."
    p.font.size = Pt(14)
    p = tf4.add_paragraph()
    p.text = "• Lysosomes & Centrosomes: Predominantly active in animal cells for autophagy and mitotic spindle formation."
    p.font.size = Pt(14)

    prs.save(pptx_path)
    print(f"Generated PPTX: {pptx_path} ({os.path.getsize(pptx_path)} bytes)")

def generate_txt_history():
    txt_path = os.path.join(FIXTURES_DIR, "industrial_revolution.txt")
    content = """The Industrial Revolution: Transformations in Energy, Labor, and Global Society (1760-1850)

Overview & Origins in Great Britain:
The Industrial Revolution was the historic transition from agrarian, handcrafted economies to industrial systems dominated by mechanized manufacture, fossil-fuel power, and factory production. Originating in Great Britain during the late 18th century, several converging catalysts facilitated this transformation:
1. Vast domestic coal and iron ore deposits situated close to navigable rivers and seaports.
2. Agricultural Revolution innovations (enclosure acts, crop rotation, selective breeding) releasing surplus agrarian labor.
3. Access to global colonial markets and raw materials (such as raw cotton from the Americas and India).
4. Stable legal and financial institutions protecting patent property rights and enabling capital investment.

Key Innovations and Mechanization:
- Textile Industry: The Flying Shuttle (John Kay, 1733), Spinning Jenny (James Hargreaves, 1764), and Water Frame (Richard Arkwright, 1769) multiplied yarn and cloth production exponentially.
- Steam Power: James Watt's separate condenser steam engine (patented 1769) transformed thermal energy from burning coal into continuous rotary mechanical power, liberating factories from dependence on river water wheels.
- Metallurgy: Abraham Darby and Henry Cort's puddling process allowed coke-smelted pig iron to be refined into durable wrought iron for bridges, rails, and steam boilers.
- Transportation Revolution: George Stephenson's Rocket (1829) and the expansion of railways connected industrial hubs like Manchester and Liverpool, accelerating raw material intake and manufactured goods distribution.

Social and Economic Consequences:
- Rapid Urbanization: Populations migrated from rural villages to industrial cities (e.g., Manchester, Birmingham, Leeds), leading to dense tenement housing, sanitation crises, and cholera outbreaks.
- The Factory System: Work shifted from seasonal, artisanal rhythms to strict, clock-disciplined 12-to-16 hour shifts under dangerous mechanical conditions.
- Child and Female Labor: Widespread employment of women and children in textile mills and coal mines due to lower wage costs and dexterous fingers for machinery.
- Social Legislation and Resistance: The Luddite movement (1811-1816) destroyed automated looms in protest of wage degradation. In response to investigative parliamentary commissions, Britain passed the Factory Act of 1833 (regulating child working hours) and Mines Act of 1842.
- Rise of Industrial Capitalism and Socialist Critique: Karl Marx and Friedrich Engels observed Manchester textile conditions, formulating historical materialism and 'The Communist Manifesto' (1848) in response to bourgeois capital accumulation and proletarian exploitation.
"""
    with open(txt_path, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"Generated TXT: {txt_path} ({os.path.getsize(txt_path)} bytes)")

def generate_edge_fixtures():
    # Empty PDF (0 bytes)
    empty_pdf = os.path.join(FIXTURES_DIR, "empty_document.pdf")
    with open(empty_pdf, "wb") as f:
        pass
    print(f"Generated empty file: {empty_pdf} (0 bytes)")

    # Corrupted DOCX (invalid magic bytes)
    corrupt_docx = os.path.join(FIXTURES_DIR, "corrupted_format.docx")
    with open(corrupt_docx, "wb") as f:
        f.write(b"CORRUPTED_NOT_A_ZIP_OR_DOCX_FILE_HEADER_1234567890\x00\xff\xfe\x00\x00\x00")
    print(f"Generated corrupt file: {corrupt_docx} ({os.path.getsize(corrupt_docx)} bytes)")

    # Large syllabus text
    large_txt = os.path.join(FIXTURES_DIR, "large_syllabus.txt")
    with open(large_txt, "w", encoding="utf-8") as f:
        f.write("# Comprehensive Multi-Discipline STEM Curriculum\n\n")
        for chapter in range(1, 51):
            f.write(f"## Chapter {chapter}: Advanced Applied Principles in Science and Technology\n")
            f.write(f"Section {chapter}.1: Fundamental Theorems, Empirical Observations, and Theoretical Models.\n")
            f.write("In this section we explore rigorous mathematical derivations, state variable equations, and computational heuristics. " * 8 + "\n")
            f.write(f"Section {chapter}.2: Practical Laboratory Implementations, Case Studies, and Edge Boundary Conditions.\n")
            f.write("Students analyze experimental data, error bounds, algorithmic asymptotic complexity, and systemic safety constraints. " * 8 + "\n\n")
    print(f"Generated large syllabus: {large_txt} ({os.path.getsize(large_txt)} bytes)")

if __name__ == "__main__":
    generate_pdf_calculus()
    generate_docx_bst()
    generate_pptx_biology()
    generate_txt_history()
    generate_edge_fixtures()
    print("All fixtures generated successfully!")
